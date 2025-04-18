from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Single slide
slide_layout = prs.slide_layouts[1]  # Title and content
slide = prs.slides.add_slide(slide_layout)

# Title
title = slide.shapes.title
title.text = "Persona Weight Calculation Logic"

# Left half: Bullet points
left = Inches(0.5)
top = Inches(1)
width = Inches(6)
height = Inches(5)
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.text = (
    "Base Weight (Pick Plan/Plan Comparison):\n"
    "- Base weight starts at 0.\n"
    "- If user picked a plan (`plan_id` exists):\n"
    "  - Use the plan’s value for the persona (e.g., `csnp` value).\n"
    "  - Limit to 0.7 for `csnp`, 0.5 for others.\n"
    "  - Multiply by:\n"
    "    - 7.0 for `csnp` if plan is special (`csnp_type = 'Y'`), else 3.0.\n"
    "    - 1.8 for `dsnp` if special (`dsnp_type = 'Y'`), else 1.2.\n"
    "    - 1.0 for other personas.\n"
    "- If no plan but user compared plans (`compared_plan_ids` exists):\n"
    "  - Average the persona’s value across compared plans.\n"
    "  - Limit to 0.7 for `csnp`, 0.5 for others.\n"
    "  - Multiply by a value based on how many compared plans are special.\n"
    "- If neither, base weight stays 0.\n\n"
    "Behavior Weight:\n"
    "- Check user actions:\n"
    "  - Queries (e.g., searches for `csnp`).\n"
    "  - Filters (e.g., filtering by `csnp`).\n"
    "  - Page views (up to 3 pages).\n"
    "  - Clicks (e.g., `pro_click_count` for `doctor`).\n"
    "- Multiply actions by:\n"
    "  - Queries: 2.5 for `csnp`, 0.8 for others.\n"
    "  - Filters: 2.3 for `csnp`, 0.7 for others.\n"
    "  - Page views: 0.15 for all.\n"
    "  - Clicks: 0.4 for `doctor`, 0.3 for `drug`, 0 for others.\n"
    "- Add extra points if:\n"
    "  - User used filters and clicked: +0.8.\n"
    "  - User used filters or clicked: +0.4.\n"
    "- Add special points for each persona:\n"
    "  - Doctor: +0.5 if clicks ≥ 1.5, +0.25 if clicks ≥ 0.5.\n"
    "  - Drug: +0.5 if clicks ≥ 5, +0.25 if clicks ≥ 2.\n"
    "  - Dental: +0.7 if 2+ actions, +0.4 if 1+ action, +0.6 if high quality, +0.4 if dental interaction > 0.\n"
    "  - Vision: +0.6 if 1+ action, +0.6 if high quality, +0.4 if vision interaction > 0.\n"
    "  - CSNP: +1.2 if 2+ actions, +0.8 if 1+ action, +1.2 if csnp interaction > 0, +1.0 if special plan, +0.8 if drug interaction > 0, +0.6 if doctor interaction > 0, +1.5 if high quality.\n"
    "  - OTC/Transportation: +0.5 if 1+ action, +0.5 if high quality.\n"
    "  - DSNP: No extra points beyond actions.\n\n"
    "Combine Scores:\n"
    "- Add base weight and behavior weight to get adjusted weight.\n\n"
    "Notes:\n"
    "- If the row’s persona matches (e.g., `csnp` for `csnp` weight), calculate simpler weights for other personas (e.g., `doctor`, `drug`), find the highest other weight, and ensure adjusted weight is at least highest other weight + 0.2.\n"
    "- Cap weight at 3.5 for `csnp`, 1.2 for others.\n"
    "- Return the final weight (e.g., `w_csnp` = 2.5)."
)
text_frame.word_wrap = True
for paragraph in text_frame.paragraphs:
    paragraph.font.size = Pt(10)  # Small font to fit
    paragraph.alignment = PP_ALIGN.LEFT

# Right half: Placeholder for flowchart
left = Inches(7)
top = Inches(1)
width = Inches(3)
height = Inches(5)
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.text = (
    "Flowchart (Add Manually):\n"
    "- User Data\n"
    "- Decision: Plan ID or Compared Plans?\n"
    "  - Yes: Base Weight = Plan value × multiplier\n"
    "  - No: Base Weight = 0\n"
    "- Behavior Weight: Actions × coeffs + bonuses\n"
    "- Adjusted Weight = Base + Behavior\n"
    "- Decision: Persona Matches?\n"
    "  - Yes: Ensure ≥ max other + 0.2\n"
    "  - No: Skip\n"
    "- Cap: 3.5 (csnp) or 1.2 (others)\n"
    "- Final Weight (e.g., w_csnp)"
)
text_frame.word_wrap = True
for paragraph in text_frame.paragraphs:
    paragraph.font.size = Pt(10)
    paragraph.alignment = PP_ALIGN.LEFT

# Save the presentation
prs.save("Persona_Weight_Calculation_Single_Slide.pptx")
print("PPT file saved as Persona_Weight_Calculation_Single_Slide.pptx")
